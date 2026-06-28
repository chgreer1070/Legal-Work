"""
Email and document to PDF converter.

Handles:
  - .msg  → PDF (email body + metadata header)
  - .docx → PDF via mammoth + weasyprint
  - .doc  → PDF via LibreOffice (best effort) or text extraction
  - .pptx → PDF via python-pptx + reportlab
  - .ppt  → PDF via LibreOffice (best effort)
  - .pdf  → copied as-is with metadata preserved
"""

import html
import os
import logging
import re
import shutil
import subprocess
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

import extract_msg
import mammoth
import weasyprint
from pptx import Presentation
from pptx.util import Pt
from pypdf import PdfReader, PdfWriter
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    HRFlowable,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)


# ── helpers ──────────────────────────────────────────────────────────────────

def _safe_str(value) -> str:
    if value is None:
        return ""
    return str(value).strip()


def _sanitize_filename(name: str, max_len: int = 80) -> str:
    """Strip characters invalid in file-system names."""
    name = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "_", name)
    name = name.strip(". ")
    return name[:max_len] or "untitled"


def _stamp_pdf_metadata(src: Path, dest: Path, meta: dict):
    """
    Copy src PDF to dest and write metadata into the PDF info dictionary.
    Falls back to plain copy on failure.
    """
    try:
        reader = PdfReader(str(src))
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)

        pdf_meta: dict[str, str] = {}
        if meta.get("title"):
            pdf_meta["/Title"] = meta["title"]
        if meta.get("author"):
            pdf_meta["/Author"] = meta["author"]
        if meta.get("subject"):
            pdf_meta["/Subject"] = meta["subject"]
        if meta.get("creator"):
            pdf_meta["/Creator"] = meta["creator"]
        if meta.get("date"):
            pdf_meta["/CreationDate"] = meta["date"]

        # Encode additional email fields into /Keywords
        kw_parts = []
        for k in ("from", "to", "cc", "date_str"):
            if meta.get(k):
                kw_parts.append(f"{k.upper()}: {meta[k]}")
        if kw_parts:
            pdf_meta["/Keywords"] = " | ".join(kw_parts)

        if pdf_meta:
            writer.add_metadata(pdf_meta)

        with open(dest, "wb") as f:
            writer.write(f)
    except Exception as e:
        logger.warning("PDF metadata stamping failed: %s", e)
        shutil.copy2(src, dest)


# ── email → PDF ───────────────────────────────────────────────────────────────

def _html_to_plain(html_body: str) -> str:
    """Lightweight HTML → plain text for email bodies."""
    text = re.sub(r"<br\s*/?>", "\n", html_body, flags=re.IGNORECASE)
    text = re.sub(r"</p>", "\n\n", text, flags=re.IGNORECASE)
    text = re.sub(r"</div>", "\n", text, flags=re.IGNORECASE)
    text = re.sub(r"<[^>]+>", "", text)
    text = html.unescape(text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _wrap_for_paragraph(text: str, style) -> list:
    """Split text into Paragraph flowables, preserving blank-line paragraphs."""
    flowables = []
    for block in text.split("\n\n"):
        block = block.strip()
        if not block:
            flowables.append(Spacer(1, 6))
            continue
        block = block.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        block = block.replace("\n", "<br/>")
        flowables.append(Paragraph(block, style))
        flowables.append(Spacer(1, 4))
    return flowables


def msg_to_pdf(msg_path: Path, output_dir: Path) -> dict:
    """
    Parse a .msg file and produce:
      - <output_dir>/<safe_subject>_email.pdf  (cover page + body)
      - converted PDFs for each attachment under <output_dir>/attachments/

    Returns a dict with keys: email_pdf, attachments (list of Path), metadata (dict)
    """
    msg = extract_msg.openMsg(str(msg_path))

    # ── gather metadata ──────────────────────────────────────────────────────
    subject   = _safe_str(msg.subject) or "No Subject"
    sender    = _safe_str(msg.sender)
    to        = _safe_str(msg.to)
    cc        = _safe_str(msg.cc) if hasattr(msg, "cc") else ""
    date_obj  = msg.date
    date_str  = date_obj.strftime("%B %d, %Y %I:%M %p") if date_obj else "Unknown Date"
    date_pdf  = date_obj.strftime("D:%Y%m%d%H%M%S") if date_obj else ""

    meta = {
        "title":    subject,
        "author":   sender,
        "subject":  subject,
        "creator":  "Outlook-to-PDF Converter",
        "date":     date_pdf,
        "from":     sender,
        "to":       to,
        "cc":       cc,
        "date_str": date_str,
    }

    # ── build email PDF via reportlab ────────────────────────────────────────
    safe_subject   = _sanitize_filename(subject)
    email_pdf_path = output_dir / f"{safe_subject}_email.pdf"

    styles     = getSampleStyleSheet()
    label_style = ParagraphStyle(
        "Label",
        parent=styles["Normal"],
        fontSize=9, leading=14,
        textColor=colors.HexColor("#555555"),
        fontName="Helvetica-Bold",
    )
    header_style = ParagraphStyle(
        "EmailHeader",
        parent=styles["Normal"],
        fontSize=9, leading=14,
        textColor=colors.HexColor("#333333"),
    )
    title_style = ParagraphStyle(
        "EmailTitle",
        parent=styles["Normal"],
        fontSize=14, leading=18,
        fontName="Helvetica-Bold",
        textColor=colors.HexColor("#1a1a2e"),
        spaceAfter=6,
    )
    body_style = ParagraphStyle(
        "EmailBody",
        parent=styles["Normal"],
        fontSize=10, leading=15,
        textColor=colors.HexColor("#222222"),
    )

    doc = SimpleDocTemplate(
        str(email_pdf_path),
        pagesize=letter,
        rightMargin=0.75 * inch,
        leftMargin=0.75 * inch,
        topMargin=0.75 * inch,
        bottomMargin=0.75 * inch,
    )

    story = []
    story.append(Paragraph(html.escape(subject), title_style))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#cccccc")))
    story.append(Spacer(1, 8))

    meta_rows = [["From:", sender], ["To:", to]]
    if cc:
        meta_rows.append(["CC:", cc])
    meta_rows.append(["Date:", date_str])

    table = Table(
        [[Paragraph(r[0], label_style), Paragraph(html.escape(r[1]), header_style)]
         for r in meta_rows],
        colWidths=[0.75 * inch, 6.0 * inch],
    )
    table.setStyle(TableStyle([
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ("TOPPADDING", (0, 0), (-1, -1), 2),
    ]))
    story.append(table)
    story.append(Spacer(1, 12))
    story.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#eeeeee")))
    story.append(Spacer(1, 12))

    body_text = ""
    if msg.htmlBody:
        raw = msg.htmlBody
        if isinstance(raw, bytes):
            raw = raw.decode("utf-8", errors="replace")
        body_text = _html_to_plain(raw)
    elif msg.body:
        body_text = _safe_str(msg.body)

    if body_text:
        story.extend(_wrap_for_paragraph(body_text, body_style))
    else:
        story.append(Paragraph("<i>[No message body]</i>", body_style))

    doc.build(story)

    # Stamp metadata
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        tmp_path = Path(tmp.name)
    shutil.move(str(email_pdf_path), tmp_path)
    _stamp_pdf_metadata(tmp_path, email_pdf_path, meta)
    tmp_path.unlink(missing_ok=True)

    # ── attachments ──────────────────────────────────────────────────────────
    att_dir = output_dir / "attachments"
    att_dir.mkdir(exist_ok=True)

    converted_attachments: list[Path] = []
    for attachment in msg.attachments:
        att_name = _safe_str(
            getattr(attachment, "longFilename", "")
            or getattr(attachment, "shortFilename", "")
            or "attachment"
        )
        att_data = attachment.data
        if att_data is None:
            continue

        ext_lower = Path(att_name).suffix.lower()
        safe_att  = _sanitize_filename(Path(att_name).stem)
        # Write to a sanitized basename, never the raw attacker-controlled
        # att_name (which may contain '..' or an absolute path). The containment
        # check is defense in depth in case sanitization is ever weakened.
        raw_path  = att_dir / f"{safe_att}{ext_lower}"
        if att_dir.resolve() not in raw_path.resolve().parents:
            raise ValueError(f"Unsafe attachment filename: {att_name!r}")

        with open(raw_path, "wb") as f:
            f.write(att_data)

        att_meta = {
            **meta,
            "title":   f"{subject} – {att_name}",
            "subject": f"Attachment: {att_name}",
        }

        dest = att_dir / f"{safe_att}.pdf"

        if ext_lower == ".pdf":
            _stamp_pdf_metadata(raw_path, dest, att_meta)
            raw_path.unlink(missing_ok=True)
            converted_attachments.append(dest)

        elif ext_lower in (".doc", ".docx", ".ppt", ".pptx"):
            pdf_path = convert_office_to_pdf(raw_path, att_dir, att_meta)
            if pdf_path and pdf_path != dest:
                shutil.move(str(pdf_path), dest)
            elif pdf_path:
                dest = pdf_path
            raw_path.unlink(missing_ok=True)
            converted_attachments.append(dest if dest.exists() else raw_path)

        else:
            # Unsupported attachment type – keep original
            converted_attachments.append(raw_path)

    msg.close()

    return {
        "email_pdf":   email_pdf_path,
        "attachments": converted_attachments,
        "metadata":    meta,
    }


# ── office document → PDF ─────────────────────────────────────────────────────

def _docx_to_pdf(src: Path, output_dir: Path, meta: dict) -> Optional[Path]:
    """Convert .docx to PDF using mammoth (HTML) + weasyprint."""
    try:
        with open(src, "rb") as f:
            result = mammoth.convert_to_html(f)
        html_body = result.value

        html_doc = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"/>
<style>
  body {{
    font-family: Georgia, 'Times New Roman', serif;
    font-size: 12pt;
    line-height: 1.6;
    margin: 2.5cm 2.5cm 2.5cm 2.5cm;
    color: #111;
  }}
  h1, h2, h3, h4, h5, h6 {{
    font-family: Arial, Helvetica, sans-serif;
    margin-top: 1em;
    margin-bottom: .4em;
  }}
  h1 {{ font-size: 18pt; }}
  h2 {{ font-size: 15pt; }}
  h3 {{ font-size: 13pt; }}
  p {{ margin: 0 0 .6em 0; }}
  table {{ border-collapse: collapse; width: 100%; margin: .8em 0; }}
  td, th {{ border: 1px solid #aaa; padding: 4pt 6pt; }}
  th {{ background: #f0f0f0; font-weight: bold; }}
  ul, ol {{ margin: .4em 0 .6em 1.5em; }}
  img {{ max-width: 100%; }}
  @page {{ margin: 0; }}
</style>
</head>
<body>{html_body}</body>
</html>"""

        dest = output_dir / (src.stem + ".pdf")
        weasyprint.HTML(string=html_doc, base_url=str(src.parent)).write_pdf(str(dest))

        # Stamp metadata
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_path = Path(tmp.name)
        shutil.move(str(dest), tmp_path)
        _stamp_pdf_metadata(tmp_path, dest, meta)
        tmp_path.unlink(missing_ok=True)

        return dest
    except Exception as e:
        logger.warning("Document conversion failed: %s", e)
        return None


def _pptx_to_pdf(src: Path, output_dir: Path, meta: dict) -> Optional[Path]:
    """Convert .pptx to PDF using python-pptx text extraction + reportlab."""
    try:
        prs = Presentation(str(src))

        styles   = getSampleStyleSheet()
        slide_title_style = ParagraphStyle(
            "SlideTitle",
            parent=styles["Normal"],
            fontSize=16, leading=20,
            fontName="Helvetica-Bold",
            textColor=colors.HexColor("#1a1a2e"),
            spaceAfter=8,
        )
        slide_num_style = ParagraphStyle(
            "SlideNum",
            parent=styles["Normal"],
            fontSize=8, leading=10,
            textColor=colors.HexColor("#888888"),
            spaceAfter=4,
        )
        body_para_style = ParagraphStyle(
            "SlideBody",
            parent=styles["Normal"],
            fontSize=11, leading=16,
            textColor=colors.HexColor("#222222"),
            leftIndent=12,
        )

        dest = output_dir / (src.stem + ".pdf")
        doc = SimpleDocTemplate(
            str(dest),
            pagesize=letter,
            rightMargin=0.75 * inch,
            leftMargin=0.75 * inch,
            topMargin=0.75 * inch,
            bottomMargin=0.75 * inch,
        )

        story = []
        total_slides = len(prs.slides)

        for slide_idx, slide in enumerate(prs.slides, start=1):
            story.append(Paragraph(f"Slide {slide_idx} / {total_slides}", slide_num_style))
            story.append(HRFlowable(width="100%", thickness=0.5,
                                    color=colors.HexColor("#cccccc")))
            story.append(Spacer(1, 4))

            slide_title = ""
            body_texts: list[str] = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                # Heuristic: first non-empty text in a title placeholder is the title
                if hasattr(shape, "placeholder_format") and \
                   shape.placeholder_format is not None and \
                   shape.placeholder_format.idx == 0 and not slide_title:
                    slide_title = text
                else:
                    body_texts.append(text)

            if slide_title:
                escaped = slide_title.replace("&", "&amp;").replace("<", "&lt;")
                story.append(Paragraph(escaped, slide_title_style))

            for bt in body_texts:
                for line in bt.split("\n"):
                    line = line.strip()
                    if not line:
                        story.append(Spacer(1, 4))
                        continue
                    line_esc = line.replace("&", "&amp;").replace("<", "&lt;")
                    story.append(Paragraph(line_esc, body_para_style))

            story.append(Spacer(1, 20))

            # Page break between slides (except last)
            if slide_idx < total_slides:
                from reportlab.platypus import PageBreak
                story.append(PageBreak())

        doc.build(story)

        # Stamp metadata
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_path = Path(tmp.name)
        shutil.move(str(dest), tmp_path)
        _stamp_pdf_metadata(tmp_path, dest, meta)
        tmp_path.unlink(missing_ok=True)

        return dest
    except Exception as e:
        logger.warning("PPTX conversion failed: %s", e)
        return None


def _libreoffice_to_pdf(src: Path, output_dir: Path) -> Optional[Path]:
    """
    Attempt LibreOffice headless conversion.
    Returns Path to the output PDF or None on failure.
    """
    try:
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pdf",
             "--outdir", str(output_dir), str(src)],
            capture_output=True, text=True, timeout=120,
        )
        candidate = output_dir / (src.stem + ".pdf")
        return candidate if candidate.exists() else None
    except (subprocess.TimeoutExpired, FileNotFoundError):
        return None


def convert_office_to_pdf(
    src: Path,
    output_dir: Path,
    meta: Optional[dict] = None,
) -> Optional[Path]:
    """
    Convert a .doc/.docx/.ppt/.pptx file to PDF.

    Strategy:
      .docx → mammoth + weasyprint (preferred), fallback to LibreOffice
      .doc  → LibreOffice (best effort)
      .pptx → python-pptx + reportlab (preferred), fallback to LibreOffice
      .ppt  → LibreOffice (best effort)
    """
    if meta is None:
        meta = {
            "title":   src.stem,
            "creator": "Outlook-to-PDF Converter",
            "date":    datetime.now().strftime("D:%Y%m%d%H%M%S"),
        }

    ext_lower = src.suffix.lower()

    if ext_lower == ".docx":
        pdf = _docx_to_pdf(src, output_dir, meta)
        if pdf:
            return pdf
        # Fallback
        return _libreoffice_to_pdf(src, output_dir)

    if ext_lower == ".doc":
        return _libreoffice_to_pdf(src, output_dir)

    if ext_lower == ".pptx":
        pdf = _pptx_to_pdf(src, output_dir, meta)
        if pdf:
            return pdf
        return _libreoffice_to_pdf(src, output_dir)

    if ext_lower == ".ppt":
        return _libreoffice_to_pdf(src, output_dir)

    return None


# ── standalone file converter ─────────────────────────────────────────────────

def convert_standalone(src: Path, output_dir: Path) -> Optional[Path]:
    """
    Convert a standalone file (no email context) to PDF.
    Supports .pdf (pass-through), .doc/.docx, .ppt/.pptx.
    """
    ext_lower = src.suffix.lower()
    safe_stem = _sanitize_filename(src.stem)
    dest      = output_dir / f"{safe_stem}.pdf"

    meta = {
        "title":   src.stem,
        "creator": "Outlook-to-PDF Converter",
        "date":    datetime.now().strftime("D:%Y%m%d%H%M%S"),
    }

    if ext_lower == ".pdf":
        _stamp_pdf_metadata(src, dest, meta)
        return dest

    if ext_lower in (".doc", ".docx", ".ppt", ".pptx"):
        pdf = convert_office_to_pdf(src, output_dir, meta)
        if pdf and pdf != dest:
            shutil.move(str(pdf), dest)
            return dest
        return pdf

    return None
